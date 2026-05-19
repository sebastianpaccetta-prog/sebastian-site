"""
Ingestion script.

Reads every PDF and PPTX in ``DATA_DIR``, chunks them, embeds them with
OpenAI, and stores them in a persistent Chroma collection.

Run with:

    python -m app.ingest

It is idempotent at the file level: re-running deletes existing chunks for
each filename before re-inserting them. Drop new files into ``data/source``
and re-run to update the index.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path
from typing import Iterable

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader
from langchain_openai import OpenAIEmbeddings
from langchain_chroma import Chroma
from pptx import Presentation

from .config import settings

SUPPORTED_EXT = {".pdf", ".pptx"}


# -----------------------------
# Loaders
# -----------------------------

def load_pdf(path: Path) -> list[Document]:
    """Load a PDF as one Document per page, preserving page numbers."""
    loader = PyPDFLoader(str(path))
    docs = loader.load()
    for d in docs:
        d.metadata["source"] = path.name
        d.metadata["doc_type"] = "pdf"
        # PyPDFLoader uses 0-indexed pages; expose as 1-indexed for humans.
        if "page" in d.metadata:
            d.metadata["page"] = int(d.metadata["page"]) + 1
    return docs


def load_pptx(path: Path) -> list[Document]:
    """
    Load a PPTX as one Document per slide.

    Captures: slide title, body shape text, speaker notes. Image OCR is out
    of scope here; if your career deck has text baked into images, run an
    OCR pass first or replace those slides with editable text.
    """
    prs = Presentation(str(path))
    docs: list[Document] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = "\n".join(
                run.text
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text
            ).strip()
            if not text:
                continue
            if shape == slide.shapes.title and not title:
                title = text
            else:
                parts.append(text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        body = "\n".join(parts).strip()
        if not (title or body or notes):
            continue

        full = ""
        if title:
            full += f"# {title}\n\n"
        if body:
            full += body + "\n\n"
        if notes:
            full += f"[Speaker notes]\n{notes}"

        docs.append(
            Document(
                page_content=full.strip(),
                metadata={
                    "source": path.name,
                    "doc_type": "pptx",
                    "page": idx,  # treat slide # as page
                    "slide_title": title,
                },
            )
        )
    return docs


def load_one(path: Path) -> list[Document]:
    if path.suffix.lower() == ".pdf":
        return load_pdf(path)
    if path.suffix.lower() == ".pptx":
        return load_pptx(path)
    raise ValueError(f"Unsupported file: {path}")


# -----------------------------
# Chunking
# -----------------------------

def chunk_documents(docs: Iterable[Document]) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=settings.chunk_size,
        chunk_overlap=settings.chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    return splitter.split_documents(list(docs))


# -----------------------------
# Vector store
# -----------------------------

def get_vectorstore() -> Chroma:
    embeddings = OpenAIEmbeddings(
        model=settings.embed_model,
        api_key=settings.openai_api_key,
    )
    return Chroma(
        collection_name="sebastian_docs",
        embedding_function=embeddings,
        persist_directory=settings.chroma_dir,
    )


def delete_existing(vs: Chroma, source: str) -> None:
    """Delete any existing chunks belonging to ``source`` so re-ingestion is idempotent."""
    try:
        vs.delete(where={"source": source})
    except Exception as e:
        print(f"  (warn) could not delete existing chunks for {source}: {e}")


# -----------------------------
# Main
# -----------------------------

def main() -> int:
    if not settings.openai_api_key:
        print("ERROR: OPENAI_API_KEY is not set. Copy .env.example -> .env and fill it in.")
        return 1

    data_dir = Path(settings.data_dir)
    if not data_dir.exists():
        print(f"ERROR: data dir {data_dir} does not exist.")
        return 1

    files = sorted(
        p for p in data_dir.iterdir()
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXT
    )

    if not files:
        print(f"No PDF or PPTX files found in {data_dir}.")
        print("Drop your resume, career deck, and 8 project PDFs in there.")
        return 0

    Path(settings.chroma_dir).mkdir(parents=True, exist_ok=True)

    vs = get_vectorstore()
    total_chunks = 0

    for path in files:
        print(f"→ {path.name}")
        delete_existing(vs, path.name)
        try:
            raw = load_one(path)
        except Exception as e:
            print(f"  (skip) failed to load: {e}")
            continue
        chunks = chunk_documents(raw)
        if not chunks:
            print("  (skip) produced no chunks")
            continue
        vs.add_documents(chunks)
        total_chunks += len(chunks)
        print(f"  + {len(chunks)} chunks ingested")

    print(f"\nDone. {total_chunks} chunks across {len(files)} files.")
    print(f"Persisted to: {settings.chroma_dir}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
